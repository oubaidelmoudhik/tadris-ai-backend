"""
PPTX Service - Extract metadata and text from PowerPoint files.
Migrated from Flask's preprocess_data.py
"""
import os
import re
from pptx import Presentation
from typing import Dict, Any, Optional


# Directory for storing lesson files (fallback if not provided)
LESSONS_DIR = "lessons"


def extract_metadata_from_filename(filename: str) -> Dict[str, Any]:
    """
    Extract metadata from filename.
    Supports formats:
    1. FR_N5_P1_SEM1_S3_V2.pptx (Short code)
    2. Français_Niv5_Parcour1_Palier3_Séance1.pptx (Long format)
    
    Args:
        filename: The PPTX filename to parse
        
    Returns:
        Dictionary with title, subject, level, period, week, session
    """
    name, _ = os.path.splitext(filename)
    parts = name.split("_")
    
    metadata = {
        "title": name.replace("_", " "),
        "subject": "français",  # Default
        "level": "5",
        "period": "",
        "week": "",
        "session": ""
    }

    # Try Short Format: FR_N5_P1_SEM1_S3...
    if len(parts) >= 5 and len(parts[0]) <= 4:
        # Subject
        if parts[0].upper() == "FR":
            metadata["subject"] = "français"
        elif parts[0].upper() == "MATH":
            metadata["subject"] = "mathématiques"
        elif parts[0].upper() == "AR":
            metadata["subject"] = "langue arabe"
        
        # Level
        if len(parts) > 1 and parts[1].startswith("N"):
            metadata["level"] = parts[1][1:]
        
        # Period (P)
        if len(parts) > 2 and parts[2].startswith("P"):
            metadata["period"] = parts[2][1:]
        
        # Semaine (SEM) -> Week
        if len(parts) > 3 and parts[3].startswith("SEM"):
            metadata["week"] = parts[3][3:]
        
        # Seance (S) -> Session
        if len(parts) > 4 and parts[4].startswith("S"):
            metadata["session"] = parts[4][1:]

    # Try Long Format: Français_Niv5_Parcour1_Palier3_Séance1
    else:
        for part in parts:
            lower = part.lower()
            if "français" in lower or "francais" in lower:
                metadata["subject"] = "français"
            elif "math" in lower:
                metadata["subject"] = "mathématiques"
            elif "arabe" in lower:
                metadata["subject"] = "langue arabe"
            
            if lower.startswith("niv"):
                metadata["level"] = lower.replace("niv", "")
            
            if lower.startswith("parcour"):
                # Map Parcour to Week
                metadata["week"] = re.sub(r'[^0-9]', '', part.lower().replace("parcour", ""))
            
            if lower.startswith("palier"):
                metadata["period"] = re.sub(r'[^0-9]', '', lower.replace("palier", ""))
                
            if lower.startswith("séance") or lower.startswith("seance"):
                metadata["session"] = re.sub(r'[^0-9]', '', lower.replace("séance", "").replace("seance", ""))

    return metadata


def extract_text_from_pptx(file_path: str) -> str:
    """
    Extract all text content from a PowerPoint file.
    
    Args:
        file_path: Path to the PPTX file
        
    Returns:
        Extracted text as string, or empty string on error
    """
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return ""


def extract_objective(content: str) -> str:
    """
    Try to extract objective from lesson content.
    
    Args:
        content: The lesson text content
        
    Returns:
        Extracted objective or default placeholder
    """
    lines = content.split('\n')
    for line in lines:
        if "objectif" in line.lower() or "capable de" in line.lower():
            return line.strip()
    return "......"


def get_or_create_lesson(filename: str, content: str = "", existing_lessons: list = None) -> Dict[str, Any]:
    """
    Get existing lesson or prepare metadata for new lesson.
    
    Args:
        filename: Name of the PPTX file
        content: Extracted text content
        existing_lessons: List of existing lesson dicts (from DB)
        
    Returns:
        Dictionary with metadata and content
    """
    metadata = extract_metadata_from_filename(filename)
    
    # If we have existing lessons, check for duplicate
    if existing_lessons:
        for lesson in existing_lessons:
            if lesson.get('filename') == filename:
                # Update existing
                metadata['id'] = lesson.get('id')
                metadata['content'] = content or lesson.get('content', '')
                metadata['objective'] = extract_objective(metadata['content'])
                return metadata
    
    # New lesson
    metadata['content'] = content
    metadata['objective'] = extract_objective(content)
    return metadata
